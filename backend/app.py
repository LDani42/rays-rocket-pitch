import os
import re
from flask import Flask, request, jsonify
from flask_cors import CORS
from dotenv import load_dotenv
from dotenv import find_dotenv
from openai import OpenAI  # Updated import

# Find the .env file
try:
    load_dotenv()
    print("Loaded environment from .env file")
except:
    print("No .env file found, using system environment variables")

# Get API key and print debug info (don't print the full key for security)
api_key = os.environ.get("OPENAI_API_KEY")
if api_key:
    print(f"Found API key starting with: {api_key[:4]}...")
else:
    print("WARNING: OPENAI_API_KEY environment variable not set")

# Create OpenAI client
client = OpenAI(api_key=api_key)

app = Flask(__name__)
CORS(app, resources={r"/*": {"origins": "*", "methods": ["GET", "POST"], "allow_headers": ["Content-Type", "Authorization"]}})

@app.route('/upload', methods=['POST'])
def upload_files():
    # Add debug prints
    print("Received upload request")
    
    pitch_deck = request.files.get('pitch_deck')
    audio_file = request.files.get('audio')
    
    print(f"Files received - Pitch deck: {pitch_deck is not None}, Audio: {audio_file is not None}")

    # Transcribe audio
    transcript_text = ""
    if audio_file:
        transcript_text = transcribe_audio(audio_file)

    # Extract text from pitch deck
    deck_text = ""
    if pitch_deck:
        deck_text = extract_deck_text(pitch_deck)

    # Generate analysis from GPT
    analysis = generate_analysis(transcript_text, deck_text)
    
    print("Analysis generated successfully")
    return jsonify({'analysis': analysis, 'transcript': transcript_text})

def transcribe_audio(file_obj):
    """
    Uses OpenAI Whisper API to transcribe audio
    """
    try:
        # Save temporary file for OpenAI API
        temp_file_path = "temp_audio.webm"
        file_obj.save(temp_file_path)
        
        with open(temp_file_path, "rb") as audio_file:
            transcript = client.audio.transcriptions.create(
                model="whisper-1",
                file=audio_file
            )
        
        # Clean up temp file
        os.remove(temp_file_path)
        return transcript.text
    except Exception as e:
        print(f"Transcription error: {e}")
        return f"(Error transcribing audio: {str(e)})"

def extract_deck_text(file_obj):
    """
    Extract text from PDF or PPT files
    """
    try:
        # Get file extension
        filename = file_obj.filename
        ext = filename.split('.')[-1].lower()
        
        # Save temp file
        temp_file_path = f"temp_deck.{ext}"
        file_obj.save(temp_file_path)
        
        extracted_text = ""
        
        # Handle different file types
        if ext == 'pdf':
            import PyPDF2
            with open(temp_file_path, 'rb') as f:
                pdf_reader = PyPDF2.PdfReader(f)
                for page_num in range(len(pdf_reader.pages)):
                    page = pdf_reader.pages[page_num]
                    extracted_text += page.extract_text() + "\n"
        
        elif ext in ['ppt', 'pptx']:
            from pptx import Presentation
            prs = Presentation(temp_file_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        extracted_text += shape.text + "\n"
        
        # Clean up temp file
        os.remove(temp_file_path)
        
        return extracted_text if extracted_text else "(No text extracted from deck)"
    
    except Exception as e:
        print(f"Text extraction error: {e}")
        return f"(Error extracting text: {str(e)})"

def generate_analysis(transcript, deck_text):
    # Enhanced system prompt based on rocket pitch requirements
    system_prompt = """You are an elite startup pitch advisor with extensive experience evaluating rocket pitches - 
concise 4-minute business presentations that follow a specific format. 

Your feedback is highly structured and objective, based on specific evidence from the materials. You use a detailed 
scoring system (0-100) for each section, with scores determined by specific criteria and supported by direct quotes.

When evaluating, always include direct quotes from the deck or transcript to support your findings.
Keep your analysis factual and evidence-based, not speculative. Structure your feedback with
clear section headings and bullet points for readability."""

    # Enhanced user prompt with objective scoring requirements
    user_prompt = f"""Analyze the following rocket pitch (transcript and deck) and provide a comprehensive, structured critique
based on the standard 4-minute rocket pitch format. Include verbatim quotes to support your analysis.

AUDIO TRANSCRIPT:
{transcript}

PITCH DECK CONTENT:
{deck_text}

For each section of your analysis:
1. Provide a score from 0-100 based on objective criteria
2. Include direct quotes from the materials to support your assessment
3. Make specific, actionable recommendations

Structure your analysis as follows:

1) PROBLEM FRAMING (Slide 1) - Score: X/100
- Quote relevant text from the deck/transcript that shows how the problem is articulated
- Assess if statistics and research are used effectively, with specific examples
- Evaluate how well the presenter establishes the target audience (quote exact language used)
- Provide specific recommendations based on the evidence

2) SOLUTION FRAMING (Slide 2) - Score: X/100
- Quote text that demonstrates how the solution addresses the stated problem
- Analyze how well benefits and advantages are explained, with direct examples
- Review evidence used to validate the solution (quote statistics or examples used)
- Provide recommendations based directly on gaps in the presented evidence

3) BUSINESS MODEL (Slide 3) - Score: X/100
- Quote text explaining the revenue generation approach
- Analyze market size claims with direct quotes of any figures provided
- Assess connection between business model and solution (with supporting text)
- Provide recommendations based on specific weaknesses in the presented model

4) FINANCIAL OVERVIEW (Slide 4) - Score: X/100
- Quote financial projections, transactions, costs and margins as presented
- Assess reasonableness of claims based on the specific numbers provided
- Note any gaps in financial information with specific examples
- Recommend improvements to financial presentation with direct references

5) PITCH DELIVERY (based on transcript)
(If no audio transcript provided, clearly state this and skip scoring this section)
- Quote examples of effective or ineffective delivery techniques
- Analyze pacing and time management based on transcript evidence
- Evaluate storytelling approaches with specific quoted examples
- Recommend concrete delivery improvements based on the transcript

6) EXECUTIVE SUMMARY
- Overall assessment with a total score (average of the section scores)
- List 3 specific strengths with supporting evidence
- List 3 priority improvement areas with specific, actionable recommendations
- Final verdict on how convincing the pitch is, based on the evidence

Use detailed evidence from the materials for objective scoring. Quote exact language where possible.
"""

    # Updated OpenAI API call with adjusted parameters
    response = client.chat.completions.create(
        model="gpt-4o",  # You could also use "gpt-4o" for even better analysis
        messages=[
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": user_prompt}
        ],
        max_tokens=4000,  # Increased token limit for more detailed analysis
        temperature=0.3,  # Lower temperature for more consistent, factual responses
        top_p=0.95,  # High precision in token selection
        presence_penalty=0.1,  # Slight penalty for repeated content
        frequency_penalty=0.1  # Slight penalty for repeated phrases
    )

    return response.choices[0].message.content

@app.route('/health', methods=['GET'])
def health_check():
    """Simple endpoint to verify the API is running"""
    return jsonify({"status": "ok"})

if __name__ == '__main__':
    # Debug mode for local testing only
    app.run(debug=True, host='0.0.0.0', port=5000)
